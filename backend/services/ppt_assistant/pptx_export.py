"""slide_deck → .pptx（python-pptx）：文案 + 简单版式 + 配图占位；可选整页嵌入百炼配图。"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Dict, List, Optional

from backend.services.ppt_assistant.bullets import bullet_parts, bullet_slide_hint
from backend.services.ppt_assistant.slide_text_layout import (
    effective_text_scheme,
    slide_body_long_text,
    slide_lead_text,
    slide_subtitle_text,
    TEXT_SCHEME_LONG_PROSE,
    TEXT_SCHEME_TITLE_LEAD,
    TEXT_SCHEME_TITLE_ONLY,
    TEXT_SCHEME_TITLE_SUBTITLE_LEAD,
)


def _slide_index(slide: Dict[str, Any], fallback: int) -> int:
    try:
        return int(slide.get("index", fallback))
    except (TypeError, ValueError):
        return fallback


def _blank_slide_layout(prs: Any) -> Any:
    try:
        return prs.slide_layouts[6]
    except IndexError:
        return prs.slide_layouts[-1]


def slide_deck_to_pptx_bytes(
    deck: Dict[str, Any],
    slide_images: Optional[Dict[int, str]] = None,
) -> bytes:
    """
    将 slide_deck JSON 写成 PowerPoint 字节流。

    slide_images: 页 index → 本地 PNG/JPG 绝对路径；若该页有有效图片则导出一整页满幅配图（文字已应在图中），
    否则沿用标题+要点版式（与原先一致）。
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise ImportError(
            "未安装 python-pptx，请执行: pip install python-pptx"
        ) from e

    if not isinstance(deck, dict):
        raise ValueError("slide_deck 必须是对象")
    slides = deck.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("slide_deck.slides 不能为空")

    imgs: Dict[int, str] = {}
    if slide_images:
        for k, v in slide_images.items():
            try:
                ik = int(k)
            except (TypeError, ValueError):
                continue
            p = Path(str(v))
            if p.is_file():
                imgs[ik] = str(p.resolve())

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    blank_layout = _blank_slide_layout(prs)

    deck_title = str(deck.get("deck_title") or "演示").strip() or "演示"
    valid_slides = [s for s in slides if isinstance(s, dict)]
    indexed = list(enumerate(valid_slides))
    indexed.sort(key=lambda t: _slide_index(t[1], t[0] + 1))
    sorted_slides: List[Dict[str, Any]] = [s for _, s in indexed]

    for i, s in enumerate(sorted_slides):
        kind = str(s.get("kind") or "content").strip()
        idx = _slide_index(s, i + 1)
        stitle = str(s.get("title") or "").strip() or f"第 {idx} 页"

        use_picture = idx in imgs

        if use_picture:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                imgs[idx],
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
        elif i == 0 and kind == "title":
            slide = prs.slides.add_slide(title_layout)
            if slide.shapes.title:
                slide.shapes.title.text = stitle
            sub = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if sub and hasattr(sub, "text_frame"):
                sub.text_frame.text = deck_title
        else:
            slide = prs.slides.add_slide(content_layout)
            if slide.shapes.title:
                slide.shapes.title.text = stitle
            body_ph = None
            try:
                body_ph = slide.placeholders[1]
            except IndexError:
                for ph in slide.placeholders:
                    try:
                        if ph.placeholder_format.idx == 1:
                            body_ph = ph
                            break
                    except (AttributeError, ValueError):
                        continue
            if body_ph and hasattr(body_ph, "text_frame"):
                tf = body_ph.text_frame
                tf.clear()
                bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
                scheme = effective_text_scheme(s)
                if scheme == TEXT_SCHEME_TITLE_ONLY:
                    tf.text = "（本页仅标题）"
                    tf.paragraphs[0].font.size = Pt(18)
                    tf.paragraphs[0].font.italic = True
                elif scheme == TEXT_SCHEME_TITLE_LEAD:
                    lead = slide_lead_text(s)
                    tf.text = lead or "（无短文）"
                    tf.paragraphs[0].font.size = Pt(20)
                elif scheme == TEXT_SCHEME_TITLE_SUBTITLE_LEAD:
                    sub = slide_subtitle_text(s)
                    lead = slide_lead_text(s)
                    tf.text = sub or "（无小标题）"
                    tf.paragraphs[0].font.size = Pt(18)
                    tf.paragraphs[0].font.bold = True
                    if lead:
                        p = tf.add_paragraph()
                        p.text = lead
                        p.level = 0
                        p.font.size = Pt(20)
                elif scheme == TEXT_SCHEME_LONG_PROSE:
                    prose = slide_body_long_text(s) or "（无正文）"
                    paras = [x.strip() for x in prose.split("\n\n") if x.strip()]
                    if not paras:
                        paras = [prose]
                    tf.text = paras[0]
                    tf.paragraphs[0].font.size = Pt(17)
                    for para in paras[1:]:
                        p = tf.add_paragraph()
                        p.text = para
                        p.level = 0
                        p.font.size = Pt(17)
                else:
                    first_bullet = True
                    for b in bullets:
                        t, _ = bullet_parts(b)
                        if not t:
                            continue
                        h = bullet_slide_hint(b)
                        if first_bullet:
                            tf.text = t
                            tf.paragraphs[0].level = 0
                            tf.paragraphs[0].font.size = Pt(20)
                            first_bullet = False
                        else:
                            p = tf.add_paragraph()
                            p.text = t
                            p.level = 0
                            p.font.size = Pt(20)
                        if h:
                            ph = tf.add_paragraph()
                            ph.text = h
                            ph.level = 1
                            ph.font.size = Pt(13)
                            ph.font.italic = True
                    if first_bullet:
                        tf.text = "（无要点）"
                        tf.paragraphs[0].font.size = Pt(20)

            if not use_picture:
                left = Inches(0.6)
                top = Inches(6.0)
                width = Inches(12.0)
                height = Inches(1.0)
                box = slide.shapes.add_textbox(left, top, width, height)
                btf = box.text_frame
                btf.text = "配图区（占位）：请在 PowerPoint 中插入图片或图标，替换本说明。"
                for p in btf.paragraphs:
                    p.font.size = Pt(11)
                    p.font.italic = True

        notes_parts: List[str] = []
        sn = str(s.get("speaker_notes") or "").strip()
        if sn:
            notes_parts.append(sn)
        bullets_list = s.get("bullets") if isinstance(s.get("bullets"), list) else []
        for b in bullets_list:
            t, e = bullet_parts(b)
            if e:
                notes_parts.append(f"「{t}」\n{e}")
        notes_text = "\n\n".join(notes_parts)
        if notes_text:
            try:
                ns = slide.notes_slide
                ns.notes_text_frame.text = notes_text
            except Exception:
                pass

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
